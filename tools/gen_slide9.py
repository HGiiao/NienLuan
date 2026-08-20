#!/usr/bin/env python3
"""Tạo slide 9 — Tính năng chính: Tìm kiếm & Đặt vé"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import copy

# ── Màu sắc ──
PRIMARY   = RGBColor(0x25, 0x63, 0xEB)  # xanh dương
DARK      = RGBColor(0x0F, 0x17, 0x2A)  # xanh đậm
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_BG  = RGBColor(0xF1, 0xF5, 0xF9)  # xám nhạt
ACCENT    = RGBColor(0xF5, 0x9E, 0x0B)  # vàng accent
TEXT_DARK  = RGBColor(0x1E, 0x29, 0x3B)
TEXT_MED   = RGBColor(0x47, 0x55, 0x69)
LIGHT_BLUE = RGBColor(0xDB, 0xEA, 0xFE)  # xanh nhạt cho card
GREEN      = RGBColor(0x16, 0xA3, 0x4A)

# ── Kích thước slide 16:9 ──
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

# ── Helper functions ──

def set_slide_bg(slide, color):
    """Set solid background color for slide."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_textbox(slide, left, top, width, height, text, font_size=12,
                color=TEXT_DARK, bold=False, alignment=PP_ALIGN.LEFT,
                font_name="Arial"):
    """Thêm textbox với text đơn giản."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox

def add_rounded_rect(slide, left, top, width, height, fill_color,
                     border_color=None, border_width=Pt(0)):
    """Thêm hình rounded rectangle."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = border_width
    else:
        shape.line.fill.background()
    # Bo góc nhỏ
    shape.adjustments[0] = 0.05
    return shape

def add_chevron(slide, left, top, width, height, fill_color):
    """Thêm mũi tên chevron."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.CHEVRON, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    shape.adjustments[0] = 0.25
    return shape

def add_circle(slide, left, top, size, fill_color):
    """Thêm hình tròn (dùng cho số bước)."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, left, top, size, size
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape

def set_shape_text(shape, text, font_size=11, color=WHITE, bold=True,
                   alignment=PP_ALIGN.CENTER, font_name="Arial"):
    """Set text cho shape."""
    tf = shape.text_frame
    tf.word_wrap = True
    tf.paragraphs[0].alignment = alignment
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    # Vertical center
    shape.text_frame.paragraphs[0].space_before = Pt(0)
    shape.text_frame.paragraphs[0].space_after = Pt(0)

def add_gradient_rect(slide, left, top, width, height):
    """Thêm hình chữ nhật với gradient xanh cho header area."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left, top, width, height
    )
    fill = shape.fill
    fill.gradient()
    fill.gradient_stops[0].color.rgb = PRIMARY
    fill.gradient_stops[0].position = 0.0
    fill.gradient_stops[1].color.rgb = RGBColor(0x1D, 0x4E, 0xD8)
    fill.gradient_stops[1].position = 1.0
    shape.line.fill.background()
    return shape

# ══════════════════════════════════════════════════════════════
# TẠO SLIDE
# ══════════════════════════════════════════════════════════════

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H

# Use blank layout
blank_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(blank_layout)

# ── Nền slide: trắng/xám nhạt ──
set_slide_bg(slide, WHITE)

# ── HEADER: gradient bar trên cùng ──
header_bar = add_gradient_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(1.0))

# ── Tiêu đề slide ──
add_textbox(slide, Inches(0.6), Inches(0.15), Inches(8), Inches(0.7),
            "TÍNH NĂNG CHÍNH — TÌM KIẾM & ĐẶT VÉ",
            font_size=26, color=WHITE, bold=True, font_name="Arial")

# ── Gạch chân accent ──
accent_line = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(0.85), Inches(3.5), Inches(0.04)
)
accent_line.fill.solid()
accent_line.fill.fore_color.rgb = ACCENT
accent_line.line.fill.background()

# ── Mô tả ngắn bên phải header ──
add_textbox(slide, Inches(9.0), Inches(0.2), Inches(4), Inches(0.6),
            "Luồng đặt vé 5 bước đơn giản • 1–9 hành khách • Nhiều chặng • Bảo hiểm",
            font_size=11, color=RGBColor(0xBF, 0xDB, 0xFE), bold=False,
            font_name="Arial", alignment=PP_ALIGN.RIGHT)

# ══════════════════════════════════════════════════════════════
# PHẦN 1: LUỒNG 5 BƯỚC NGANG
# ══════════════════════════════════════════════════════════════

# Vùng tiêu đề phần
add_textbox(slide, Inches(0.6), Inches(1.15), Inches(5), Inches(0.4),
            "▸ LUỒNG ĐẶT VÉ",
            font_size=14, color=PRIMARY, bold=True)

# 5 bước: tìm kiếm → chọn chuyến → nhập hành khách → thanh toán → xác nhận
steps = [
    ("1", "TÌM KIẾM", "Chọn tuyến, ngày,\nphương tiện, hạng ghế"),
    ("2", "CHỌN CHUYẾN", "Xem danh sách,\nso sánh giá, sắp xếp"),
    ("3", "NHẬP HÀNH KHÁCH", "Thông tin 1–9 người,\nbảo hiểm, liên hệ"),
    ("4", "THANH TOÁN", "5 cổng VN:\nVNPay/MoMo/ZaloPay/PayOS/VietQR"),
    ("5", "XÁC NHẬN", "Email xác nhận,\nQR code, theo dõi booking"),
]

step_colors = [PRIMARY, RGBColor(0x1D, 0x4E, 0xD8), RGBColor(0x0E, 0x74, 0x90),
               RGBColor(0x04, 0x78, 0x57), GREEN]

step_start_x = Inches(0.5)
step_y = Inches(1.65)
step_w = Inches(2.2)
step_h = Inches(1.65)
chevron_gap = Inches(0.08)
chevron_w = Inches(0.35)
chevron_h = Inches(0.5)

for i, (num, title, desc) in enumerate(steps):
    x = step_start_x + i * (step_w + chevron_gap + Inches(0.02))

    # Card nền
    card = add_rounded_rect(slide, x, step_y, step_w, step_h,
                            LIGHT_BG, border_color=step_colors[i], border_width=Pt(1.5))

    # Số bước (hình tròn nhỏ)
    circle = add_circle(slide, x + Inches(0.85), step_y + Inches(0.1),
                        Inches(0.4), step_colors[i])
    set_shape_text(circle, num, font_size=14, color=WHITE, bold=True)

    # Tên bước
    add_textbox(slide, x + Inches(0.1), step_y + Inches(0.55),
                step_w - Inches(0.2), Inches(0.35),
                title, font_size=13, color=step_colors[i], bold=True,
                alignment=PP_ALIGN.CENTER)

    # Mô tả
    add_textbox(slide, x + Inches(0.1), step_y + Inches(0.9),
                step_w - Inches(0.2), Inches(0.7),
                desc, font_size=10, color=TEXT_MED, bold=False,
                alignment=PP_ALIGN.CENTER)

    # Mũi tên chevron giữa các bước
    if i < len(steps) - 1:
        cx = x + step_w + Inches(0.01)
        cy = step_y + step_h / 2 - chevron_h / 2
        chev = add_chevron(slide, cx, cy, chevron_w, chevron_h, ACCENT)

# ══════════════════════════════════════════════════════════════
# PHẦN 2: ĐIỂM NỔI BẬT (3 cột)
# ══════════════════════════════════════════════════════════════

section2_y = Inches(3.6)
add_textbox(slide, Inches(0.6), section2_y, Inches(5), Inches(0.4),
            "▸ ĐIỂM NỔI BẬT",
            font_size=14, color=PRIMARY, bold=True)

# 3 cột nội dung
col_data = [
    {
        "icon": "✈️🚢🚌",
        "title": "3 Phương Tiện",
        "items": [
            "Máy bay: VietJet, Vietnam Airlines,\nBamboo Airways...",
            "Tàu hỏa: SE, TN, SPT,\nđường sắt Bắc–Nam",
            "Xe khách: limousine,\nghế ngồi, giường nằm",
        ],
        "color": PRIMARY,
    },
    {
        "icon": "👥",
        "title": "Đặt Vé Linh Hoạt",
        "items": [
            "1–9 hành khách / chuyến,\nphù hợp gia đình, đoàn",
            "Nhiều chặng (BookingSegments):\nHAN→SGN→DAD kết hợp",
            "Bảo hiểm 3 gói:\nCơ Bản / Cao Cấp / Toàn Diện",
        ],
        "color": RGBColor(0x0E, 0x74, 0x90),
    },
    {
        "icon": "🔍",
        "title": "Tìm Kiếm Thông Minh",
        "items": [
            "Lọc theo tuyến, ngày,\nhạng ghế, khoảng giá",
            "Sắp xếp: giá tăng dần,\ngiảm dần, thời gian bay",
            "Phân trang, xem chi tiết,\nđánh giá từ cộng đồng",
        ],
        "color": RGBColor(0x04, 0x78, 0x57),
    },
]

col_w = Inches(3.85)
col_gap = Inches(0.3)
col_start_x = Inches(0.5)
col_h = Inches(3.15)

for i, col in enumerate(col_data):
    cx = col_start_x + i * (col_w + col_gap)
    cy = section2_y + Inches(0.45)

    # Card nền
    card = add_rounded_rect(slide, cx, cy, col_w, col_h,
                            WHITE, border_color=col["color"], border_width=Pt(1.5))

    # Header card
    header_rect = add_rounded_rect(slide, cx, cy, col_w, Inches(0.55),
                                   col["color"])
    # Fix bo góc trên
    header_rect.adjustments[0] = 0.05

    # Icon + title trong header
    add_textbox(slide, cx + Inches(0.15), cy + Inches(0.05),
                col_w - Inches(0.3), Inches(0.45),
                f'{col["icon"]}  {col["title"]}',
                font_size=14, color=WHITE, bold=True,
                alignment=PP_ALIGN.LEFT)

    # Các mục
    for j, item in enumerate(col["items"]):
        iy = cy + Inches(0.7) + j * Inches(0.8)
        # Bullet marker
        bullet = add_rounded_rect(slide, cx + Inches(0.15), iy + Inches(0.08),
                                  Inches(0.08), Inches(0.08), col["color"])
        # Text
        add_textbox(slide, cx + Inches(0.35), iy - Inches(0.02),
                    col_w - Inches(0.5), Inches(0.75),
                    item, font_size=10.5, color=TEXT_DARK, bold=False)

# ══════════════════════════════════════════════════════════════
# PHẦN 3: FLOW ĐẶT VÉ CHI TIẾT (bên phải luồng 5 bước)
# ══════════════════════════════════════════════════════════════

# Thêm box nhỏ ghi chú luồng API ở góc dưới phải
note_x = Inches(9.2)
note_y = Inches(1.65)
note_w = Inches(3.8)
note_h = Inches(1.65)

note_card = add_rounded_rect(slide, note_x, note_y, note_w, note_h,
                              RGBColor(0xEF, 0xF6, 0xFF),
                              border_color=PRIMARY, border_width=Pt(1))

add_textbox(slide, note_x + Inches(0.15), note_y + Inches(0.08),
            note_w - Inches(0.3), Inches(0.3),
            "⚙️ LUỒNG API",
            font_size=12, color=PRIMARY, bold=True)

api_flow = (
    "GET /api/flights?from=HAN&to=SGN&date=...\n"
    "    ↓  Response: danh sách chuyến\n"
    "POST /api/bookings  (body: passengers[])\n"
    "    ↓  Tính giá + giảm giá + bảo hiểm\n"
    "POST /api/payments/{id}?method=vnpay\n"
    "    ↓  Redirect cổng thanh toán\n"
    "IPN/Webhook → Cập nhật Booking.Status"
)

add_textbox(slide, note_x + Inches(0.15), note_y + Inches(0.38),
            note_w - Inches(0.3), note_h - Inches(0.45),
            api_flow, font_size=9.5, color=TEXT_MED, bold=False,
            font_name="Consolas")

# ══════════════════════════════════════════════════════════════
# FOOTER — số slide
# ══════════════════════════════════════════════════════════════

# Footer bar
footer_bar = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, Inches(0), SLIDE_H - Inches(0.35),
    SLIDE_W, Inches(0.35)
)
footer_bar.fill.solid()
footer_bar.fill.fore_color.rgb = DARK
footer_bar.line.fill.background()

add_textbox(slide, Inches(0.5), SLIDE_H - Inches(0.33),
            Inches(4), Inches(0.3),
            "VÉ247 — Hệ thống đặt vé đa phương tiện thông minh",
            font_size=9, color=RGBColor(0x94, 0xA3, 0xB8), bold=False)

add_textbox(slide, SLIDE_W - Inches(2), SLIDE_H - Inches(0.33),
            Inches(1.5), Inches(0.3),
            "9 / 23",
            font_size=9, color=RGBColor(0x94, 0xA3, 0xB8), bold=False,
            alignment=PP_ALIGN.RIGHT)

# ── Lưu file ──
output_path = "docs/Ve247_Slide9_TimKiemDatVe.pptx"
prs.save(output_path)
print(f"Done: {output_path}")
